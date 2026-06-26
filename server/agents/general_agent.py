"""
GeneralAgent — General-purpose desktop assistant for non-code tasks.

Handles browse/PPT/Excel/email/calendar tasks and streams progress to the client.

WebSocket message protocol:
  chunk:      { "type": "chunk", "delta": "..." }
  tool_start: { "type": "tool_start", "tool": "<name>" }
  tool_done:  { "type": "tool_done", "tool": "<name>", "result": {...} }
  done:       { "type": "done", "model": "<name>", "latency_ms": N }
  error:      { "type": "error", "message": "..." }

Tier behaviour:
  Free       — text-only with claude-haiku-4-5; no tools exposed.
  Pro        — full tool suite (PPT, Excel, email, calendar) with claude-sonnet-4-6.
  Enterprise — same as Pro.
"""

from __future__ import annotations

import asyncio
import json
import os
import tempfile
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List, Union

import anthropic


# ---------------------------------------------------------------------------
# Tier identifiers — mirrors server/tiers.py TierId values
# ---------------------------------------------------------------------------
TIER_FREE = "free"
TIER_PRO = "pro"
TIER_ENTERPRISE = "enterprise"

TierLike = Union[str, Any]

# ---------------------------------------------------------------------------
# Model selection
# ---------------------------------------------------------------------------
_FREE_MODEL = "claude-haiku-4-5"
_PRO_MODEL = "claude-sonnet-4-6"

_MODEL_FOR_TIER: Dict[str, str] = {
    TIER_FREE: _FREE_MODEL,
    TIER_PRO: _PRO_MODEL,
    TIER_ENTERPRISE: _PRO_MODEL,
}

MAX_TOOL_ITERATIONS = 5

# ---------------------------------------------------------------------------
# System prompt
# ---------------------------------------------------------------------------
GENERAL_SYSTEM_PROMPT = (
    "You are a helpful desktop assistant integrated into OpenUI. "
    "Help users with general tasks: creating presentations, building spreadsheets, "
    "drafting emails, managing calendar events, and answering questions. "
    "When tools are available, use them to accomplish tasks rather than just describing how. "
    "After using tools, summarize the results clearly and concisely."
)

# ---------------------------------------------------------------------------
# Pro tool definitions (Anthropic JSON schema format)
# ---------------------------------------------------------------------------
_PRO_TOOLS: List[Dict] = [
    {
        "name": "create_presentation",
        "description": (
            "Create a PowerPoint presentation (.pptx) and save it to disk. "
            "Returns a download path and URL for the file."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Presentation title shown on the title slide.",
                },
                "slides": {
                    "type": "array",
                    "description": "Ordered list of content slides to generate.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide heading."},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Bullet point lines for the slide body.",
                            },
                        },
                        "required": ["title", "bullets"],
                    },
                },
            },
            "required": ["title", "slides"],
        },
    },
    {
        "name": "create_spreadsheet",
        "description": (
            "Create an Excel workbook (.xlsx) and save it to disk. "
            "Returns a download path and URL for the file."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Desired filename without extension.",
                },
                "data": {
                    "type": "array",
                    "description": "List of worksheets to create.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "sheet": {"type": "string", "description": "Sheet (tab) name."},
                            "rows": {
                                "type": "array",
                                "description": "Row data — each row is an array of cell values.",
                                "items": {"type": "array", "items": {}},
                            },
                        },
                        "required": ["sheet", "rows"],
                    },
                },
            },
            "required": ["filename", "data"],
        },
    },
    {
        "name": "draft_email",
        "description": (
            "Compose an email draft and return the formatted text. "
            "Never sends — the user must review and send it themselves."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "to": {"type": "string", "description": "Recipient address(es)."},
                "subject": {"type": "string", "description": "Email subject line."},
                "body": {
                    "type": "string",
                    "description": "Email body in plain text or markdown.",
                },
            },
            "required": ["to", "subject", "body"],
        },
    },
    {
        "name": "list_calendar_events",
        "description": "List calendar events within a given date range.",
        "input_schema": {
            "type": "object",
            "properties": {
                "date_range": {
                    "type": "string",
                    "description": (
                        "Human-readable range. Examples: 'today', 'this week', "
                        "'next 7 days', '2025-06-01 to 2025-06-30'."
                    ),
                },
            },
            "required": ["date_range"],
        },
    },
    {
        "name": "create_calendar_event",
        "description": "Create a new calendar event.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Event title/summary."},
                "start": {
                    "type": "string",
                    "description": "Start date-time in ISO 8601, e.g. '2025-06-15T14:00:00'.",
                },
                "end": {
                    "type": "string",
                    "description": "End date-time in ISO 8601.",
                },
                "description": {
                    "type": "string",
                    "description": "Optional event notes.",
                    "default": "",
                },
            },
            "required": ["title", "start", "end"],
        },
    },
]


# ---------------------------------------------------------------------------
# Agent
# ---------------------------------------------------------------------------

class GeneralAgent:
    """
    General-purpose assistant backed by Anthropic cloud models.

    Free tier runs text-only (haiku). Pro/Enterprise unlock the tool suite
    (PPT via python-pptx, Excel via openpyxl, email draft, Google Calendar).

    Usage::

        agent = GeneralAgent(tier="pro")
        await agent.stream(messages, websocket)

        # or with a TierId enum from server.tiers:
        from server.tiers import TierId
        agent = GeneralAgent(tier=TierId.PRO)
    """

    def __init__(self, tier: TierLike) -> None:
        self.tier: str = str(tier.value if hasattr(tier, "value") else tier).lower()
        self.model = _MODEL_FOR_TIER.get(self.tier, _FREE_MODEL)
        self._has_pro_tools = self.tier in (TIER_PRO, TIER_ENTERPRISE)
        self._client = anthropic.AsyncAnthropic(
            api_key=os.environ.get("ANTHROPIC_API_KEY", "")
        )

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def stream(self, messages: List[Dict], websocket: Any) -> None:
        """Stream a response (and any tool calls) over the WebSocket."""
        tools = _PRO_TOOLS if self._has_pro_tools else []
        # Strip system messages — passed separately as the `system` parameter
        api_messages = [m for m in messages if m.get("role") != "system"]
        start = time.time()
        try:
            await self._agentic_loop(api_messages, tools, websocket, start)
        except Exception as exc:
            await websocket.send_json({"type": "error", "message": str(exc)})

    # ------------------------------------------------------------------
    # Agentic loop
    # ------------------------------------------------------------------

    async def _agentic_loop(
        self,
        messages: List[Dict],
        tools: List[Dict],
        websocket: Any,
        start: float,
    ) -> None:
        """Run the streaming tool-use loop (up to MAX_TOOL_ITERATIONS rounds)."""
        for _ in range(MAX_TOOL_ITERATIONS):
            kwargs: Dict[str, Any] = {
                "model": self.model,
                "max_tokens": 4096,
                "system": GENERAL_SYSTEM_PROMPT,
                "messages": messages,
            }
            if tools:
                kwargs["tools"] = tools

            async with self._client.messages.stream(**kwargs) as stream:
                async for text in stream.text_stream:
                    await websocket.send_json({"type": "chunk", "delta": text})
                final_message = await stream.get_final_message()

            # Append assistant turn (serialise content blocks to plain dicts)
            messages.append({
                "role": "assistant",
                "content": [b.model_dump() for b in final_message.content],
            })

            if final_message.stop_reason == "end_turn":
                await websocket.send_json({
                    "type": "done",
                    "model": self.model,
                    "latency_ms": round((time.time() - start) * 1000),
                })
                return

            if final_message.stop_reason == "tool_use":
                tool_results = []
                for block in final_message.content:
                    if block.type == "tool_use":
                        await websocket.send_json({
                            "type": "tool_start",
                            "tool": block.name,
                        })
                        result = await self._execute_tool(block.name, block.input)
                        await websocket.send_json({
                            "type": "tool_done",
                            "tool": block.name,
                            "result": result,
                        })
                        tool_results.append({
                            "type": "tool_result",
                            "tool_use_id": block.id,
                            "content": json.dumps(result),
                        })
                messages.append({"role": "user", "content": tool_results})
                continue

            # Any other stop reason (max_tokens, etc.) — surface done and exit
            await websocket.send_json({
                "type": "done",
                "model": self.model,
                "latency_ms": round((time.time() - start) * 1000),
            })
            return

        # Fell out of the loop — iteration cap reached
        await websocket.send_json({
            "type": "done",
            "model": self.model,
            "latency_ms": round((time.time() - start) * 1000),
        })

    # ------------------------------------------------------------------
    # Tool dispatcher
    # ------------------------------------------------------------------

    async def _execute_tool(self, name: str, tool_input: Dict) -> Dict:
        """Dispatch a tool call by name and return the result dict."""
        if name == "create_presentation":
            return await self._create_presentation(
                tool_input.get("title", "Untitled"),
                tool_input.get("slides", []),
            )
        if name == "create_spreadsheet":
            return await self._create_spreadsheet(
                tool_input.get("filename", "spreadsheet"),
                tool_input.get("data", []),
            )
        if name == "draft_email":
            return await self._draft_email(
                tool_input.get("to", ""),
                tool_input.get("subject", ""),
                tool_input.get("body", ""),
            )
        if name == "list_calendar_events":
            return await self._list_calendar_events(tool_input.get("date_range", "today"))
        if name == "create_calendar_event":
            return await self._create_calendar_event(
                tool_input.get("title", ""),
                tool_input.get("start", ""),
                tool_input.get("end", ""),
                tool_input.get("description", ""),
            )
        return {"error": f"Unknown tool: {name}"}

    # ------------------------------------------------------------------
    # PPT tool
    # ------------------------------------------------------------------

    async def _create_presentation(self, title: str, slides: List[Dict]) -> Dict:
        """Create a .pptx file via python-pptx (blocking, run in thread)."""
        try:
            from pptx import Presentation  # noqa: F401 — verify import up-front
        except ImportError:
            return {
                "error": "python-pptx is not installed.",
                "hint": "pip install python-pptx",
            }

        def _build() -> str:
            from pptx import Presentation as Prs

            prs = Prs()

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = ""

            # Content slides
            bullet_layout = prs.slide_layouts[1]
            for slide_data in slides:
                s = prs.slides.add_slide(bullet_layout)
                s.shapes.title.text = slide_data.get("title", "")
                tf = s.placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(slide_data.get("bullets", [])):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        tf.add_paragraph().text = bullet

            out_path = Path(tempfile.gettempdir()) / f"{uuid.uuid4().hex}.pptx"
            prs.save(str(out_path))
            return str(out_path)

        path = await asyncio.to_thread(_build)
        filename = Path(path).name
        return {
            "success": True,
            "path": path,
            "filename": filename,
            "download_url": f"/files/{filename}",
            "slides_created": len(slides) + 1,  # +1 for title slide
        }

    # ------------------------------------------------------------------
    # Excel tool
    # ------------------------------------------------------------------

    async def _create_spreadsheet(self, filename: str, data: List[Dict]) -> Dict:
        """Create an .xlsx file via openpyxl (blocking, run in thread)."""
        try:
            import openpyxl  # noqa: F401
        except ImportError:
            return {
                "error": "openpyxl is not installed.",
                "hint": "pip install openpyxl",
            }

        def _build() -> str:
            import openpyxl as xl

            wb = xl.Workbook()
            wb.remove(wb.active)  # Drop the default empty sheet

            for sheet_data in data:
                name = sheet_data.get("sheet", "Sheet1")[:31]  # Excel name limit
                ws = wb.create_sheet(title=name)
                for row in sheet_data.get("rows", []):
                    ws.append(list(row))

            if not wb.worksheets:
                wb.create_sheet("Sheet1")

            safe = "".join(c for c in filename if c.isalnum() or c in "-_")[:50] or "spreadsheet"
            out_path = Path(tempfile.gettempdir()) / f"{uuid.uuid4().hex}_{safe}.xlsx"
            wb.save(str(out_path))
            return str(out_path)

        path = await asyncio.to_thread(_build)
        fname = Path(path).name
        return {
            "success": True,
            "path": path,
            "filename": fname,
            "download_url": f"/files/{fname}",
            "sheets_created": len(data),
        }

    # ------------------------------------------------------------------
    # Email draft tool
    # ------------------------------------------------------------------

    async def _draft_email(self, to: str, subject: str, body: str) -> Dict:
        """Return a formatted email draft. Never sends."""
        draft = f"To: {to}\nSubject: {subject}\n\n{body}"
        return {
            "success": True,
            "draft": draft,
            "note": "Draft only — review and send it yourself.",
        }

    # ------------------------------------------------------------------
    # Calendar tools
    # ------------------------------------------------------------------

    async def _list_calendar_events(self, date_range: str) -> Dict:
        """List events via Google Calendar API (requires GOOGLE_CALENDAR_CREDENTIALS_JSON)."""
        creds_json = os.environ.get("GOOGLE_CALENDAR_CREDENTIALS_JSON")
        if not creds_json:
            return {
                "success": False,
                "error": "Google Calendar is not configured.",
                "hint": "Set GOOGLE_CALENDAR_CREDENTIALS_JSON to your service-account JSON.",
            }
        try:
            return await asyncio.to_thread(self._gcal_list, creds_json, date_range)
        except Exception as exc:
            return {"success": False, "error": f"Calendar API error: {exc}"}

    def _gcal_list(self, creds_json: str, date_range: str) -> Dict:
        """Blocking Google Calendar events.list call."""
        from datetime import datetime, timezone, timedelta
        from google.oauth2 import service_account
        from googleapiclient.discovery import build

        creds = service_account.Credentials.from_service_account_info(
            json.loads(creds_json),
            scopes=["https://www.googleapis.com/auth/calendar.readonly"],
        )
        service = build("calendar", "v3", credentials=creds)

        now = datetime.now(timezone.utc)
        normalized = date_range.lower().strip()

        if "today" in normalized:
            time_min = now.replace(hour=0, minute=0, second=0, microsecond=0)
            time_max = now.replace(hour=23, minute=59, second=59, microsecond=0)
        elif "week" in normalized or "7 day" in normalized:
            time_min = now
            time_max = now + timedelta(days=7)
        else:
            time_min = now
            time_max = now + timedelta(days=7)

        result = (
            service.events()
            .list(
                calendarId="primary",
                timeMin=time_min.isoformat(),
                timeMax=time_max.isoformat(),
                maxResults=20,
                singleEvents=True,
                orderBy="startTime",
            )
            .execute()
        )
        events = [
            {
                "title": e.get("summary", "(No title)"),
                "start": e.get("start", {}).get("dateTime") or e.get("start", {}).get("date"),
                "end": e.get("end", {}).get("dateTime") or e.get("end", {}).get("date"),
                "description": e.get("description", ""),
            }
            for e in result.get("items", [])
        ]
        return {"success": True, "date_range": date_range, "events": events}

    async def _create_calendar_event(
        self, title: str, start: str, end: str, description: str
    ) -> Dict:
        """Create an event via Google Calendar API."""
        creds_json = os.environ.get("GOOGLE_CALENDAR_CREDENTIALS_JSON")
        if not creds_json:
            return {
                "success": False,
                "error": "Google Calendar is not configured.",
                "hint": "Set GOOGLE_CALENDAR_CREDENTIALS_JSON to your service-account JSON.",
            }
        try:
            return await asyncio.to_thread(
                self._gcal_insert, creds_json, title, start, end, description
            )
        except Exception as exc:
            return {"success": False, "error": f"Calendar API error: {exc}"}

    def _gcal_insert(
        self, creds_json: str, title: str, start: str, end: str, description: str
    ) -> Dict:
        """Blocking Google Calendar events.insert call."""
        from google.oauth2 import service_account
        from googleapiclient.discovery import build

        creds = service_account.Credentials.from_service_account_info(
            json.loads(creds_json),
            scopes=["https://www.googleapis.com/auth/calendar"],
        )
        service = build("calendar", "v3", credentials=creds)

        body = {
            "summary": title,
            "description": description,
            "start": {"dateTime": start, "timeZone": "UTC"},
            "end": {"dateTime": end, "timeZone": "UTC"},
        }
        created = service.events().insert(calendarId="primary", body=body).execute()
        return {
            "success": True,
            "event_id": created.get("id"),
            "title": created.get("summary"),
            "start": created.get("start", {}).get("dateTime"),
            "end": created.get("end", {}).get("dateTime"),
            "link": created.get("htmlLink"),
        }
